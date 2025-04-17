# from tasks.celery import celery
import pptx
from PyPDF2 import PdfReader
from extensions import cache  # Import Redis cache


# @celery.task
def extract_pdf_data(file_path):
    """Extracts text from PDF and caches results."""
    cached_result = cache.get(file_path)
    if cached_result:
        return cached_result  # Return cached data if available

    content = []
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                content.append(page.extract_text())
    except Exception as e:
        return None, f"PDF parsing failed: {str(e)}"
    
    result = "\n".join(content)
    cache.set(file_path, result, timeout=86400)  # Cache for 24 hours
    return result, None


# @celery.task
def extract_pptx_data(file_path):
    """Extracts text from PPTX and caches results."""
    cached_result = cache.get(file_path)
    if cached_result:
        return cached_result  # Return cached data if available

    content = []
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
    except Exception as e:
        return None, f"PPTX parsing failed: {str(e)}"
    
    result = "\n".join(content)
    cache.set(file_path, result, timeout=86400)  # Cache for 24 hours
    return result, None


