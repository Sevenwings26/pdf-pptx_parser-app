import pptx
from PyPDF2 import PdfReader
# "error": "PDF parsing failed: PdfFileReader is deprecated and was removed in PyPDF2 3.0.0. Use PdfReader instead."
# Extract pdf data 
def extract_pdf_data(file_path):
    content = [] # to store extracted text
    try:
        with open(file_path, 'rb') as file: # read in binary mode
            reader = PdfReader(file)
            for page in reader.pages:
                content.append(page.extract_text())
        # return content
    except Exception as e:
        return None, f"PDF parsing failed: {str(e)}"
    return "\n".join(content), None


# Extract pptx data
def extract_pptx_data(file_path):
    content = [] # to store extracted text
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
        # return content
    except Exception as e:
        return None, f"PPTx parsing failed: {str(e)}"

    return "\n".join(content), None
